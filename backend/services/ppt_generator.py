from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import os

class PPTGenerator:
    def __init__(self):
        self.output_dir = Path("temp/presentations")
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def generate_presentation(self, content_data: dict, filename: str) -> str:
        """Generate PowerPoint presentation"""
        
        try:
            print(f"📊 Generating PowerPoint...")
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Title Slide
            self._create_title_slide(prs, content_data.get('heading', 'Lesson'))
            
            # Slide 2: Definition
            if content_data.get('definition'):
                self._create_definition_slide(prs, content_data)
            
            # Slide 3+: Key Points (2-3 points per slide)
            points = content_data.get('points', [])
            self._create_points_slides(prs, content_data.get('heading', ''), points)
            
            # Last Slide: Summary
            self._create_summary_slide(prs, content_data.get('heading', ''))
            
            # Save
            output_path = self.output_dir / filename
            prs.save(str(output_path))
            
            print(f"✅ PowerPoint created: {output_path}")
            return str(output_path)
            
        except Exception as e:
            print(f"❌ PPT generation error: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _create_title_slide(self, prs, title_text):
        """Create title slide"""
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 64, 175)  # Blue
        
        # Title
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = title_text
        
        # Format
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _create_definition_slide(self, prs, content_data):
        """Create definition slide"""
        
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "📖 Definition"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = content_data['definition']
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.line_spacing = 1.5
    
    def _create_points_slides(self, prs, heading, points):
        """Create slides for key points"""
        
        # Group points (3 per slide)
        for i in range(0, len(points), 3):
            slide_points = points[i:i+3]
            
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = f"📌 {heading}"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(30, 64, 175)
            
            # Points
            left = Inches(1.5)
            top = Inches(2.5)
            width = Inches(7)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            for point in slide_points:
                point_text = point.get('text', '')
                p = tf.add_paragraph()
                p.text = f"• {point_text}"
                p.level = 0
                p.font.size = Pt(20)
                p.space_before = Pt(12)
                p.line_spacing = 1.3
    
    def _create_summary_slide(self, prs, heading):
        """Create summary slide"""
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 249, 255)
        
        # Title
        left = Inches(2)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(2)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Thank You!\n\nKeep Learning 📚"
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 64, 175)

ppt_generator = PPTGenerator()